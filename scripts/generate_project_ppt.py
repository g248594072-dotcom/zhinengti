# -*- coding: utf-8 -*-
"""生成「聊天记录质检工具」跨境电商风格美化 PPT。"""

from pathlib import Path
import struct

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "聊天记录质检工具_项目介绍.pptx"
ASSETS_DIR = Path(__file__).resolve().parent / "ppt_assets"

# ── 跨境电商主题色 ──
NAVY = RGBColor(0x0F, 0x17, 0x2A)          # 深海蓝（主色）
OCEAN = RGBColor(0x03, 0x69, 0xA1)         # 海洋蓝
SKY = RGBColor(0x38, 0xBD, 0xF8)            # 天蓝
CORAL = RGBColor(0xF9, 0x73, 0x16)         # 珊瑚橙（强调）
GOLD = RGBColor(0xF5, 0x9E, 0x0B)           # 金色
TEAL = RGBColor(0x14, 0xB8, 0xA6)           # 青绿
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = RGBColor(0xF0, 0xF9, 0xFF)        # 淡蓝页面底
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xBA, 0xE6, 0xFD)    # 浅蓝边框
LIGHT_BLUE = RGBColor(0xE0, 0xF2, 0xFE)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
RED_BORDER = RGBColor(0xFC, 0xA5, 0xA5)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
GREEN_ACCENT = RGBColor(0x05, 0x96, 0x69)
GREEN_BORDER = RGBColor(0x6E, 0xE7, 0xB7)
PLACEHOLDER_BG = RGBColor(0xF8, 0xFA, 0xFC)
PLACEHOLDER_BORDER = RGBColor(0x7D, 0xD3, 0xFC)

FOOTER_TEXT = "跨境电商 · 客服聊天质量管理"


# ── 基础绘图工具 ──

def _no_line(shape):
    shape.line.fill.background()


def _set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_circle(slide, x, y, size, color, alpha_layer=False):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    _no_line(c)
    return c


def _add_rounded_rect(slide, x, y, w, h, fill, border=None, border_pt=1.5):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if border:
        r.line.color.rgb = border
        r.line.width = Pt(border_pt)
    else:
        _no_line(r)
    return r


def _add_shadow_card(slide, x, y, w, h, fill=CARD_BG, border=CARD_BORDER):
    """卡片 + 偏移阴影层。"""
    _add_rounded_rect(slide, x + 0.04, y + 0.05, w, h, RGBColor(0xCB, 0xE7, 0xF5))
    return _add_rounded_rect(slide, x, y, w, h, fill, border, 1.5)


def _apply_light_bg(slide):
    """浅色内容页背景：渐变感 + 角落装饰圆。"""
    _set_slide_bg(slide, PAGE_BG)
    # 右上角大圆
    _add_circle(slide, 7.8, -1.2, 3.5, LIGHT_BLUE)
    _add_circle(slide, 8.5, -0.5, 2.0, RGBColor(0xDB, 0xEA, 0xFE))
    # 左下角装饰
    _add_circle(slide, -1.0, 5.8, 2.8, LIGHT_BLUE)
    _add_circle(slide, -0.3, 6.5, 1.5, RGBColor(0xDB, 0xEA, 0xFE))
    # 顶部色带
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = OCEAN
    _no_line(bar)
    # 底部细线
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(10), Inches(0.04))
    bot.fill.solid()
    bot.fill.fore_color.rgb = SKY
    _no_line(bot)


def _apply_section_bg(slide):
    """章节/封面页：统一亮色背景（与内容页一致）。"""
    _apply_light_bg(slide)
    # 底部装饰色条
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(10), Inches(0.12))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = OCEAN
    _no_line(stripe)
    stripe2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.22), Inches(10), Inches(0.05))
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = SKY
    _no_line(stripe2)


def _add_header(slide, title: str, icon: str = "🌐"):
    """带图标和底线的标题区。"""
    badge = _add_rounded_rect(slide, 0.55, 0.28, 0.55, 0.55, OCEAN)
    ip = badge.text_frame.paragraphs[0]
    ip.text = icon
    ip.font.size = Pt(18)
    ip.alignment = PP_ALIGN.CENTER

    tp = slide.shapes.add_textbox(Inches(1.25), Inches(0.28), Inches(8.0), Inches(0.65)).text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.0), Inches(8.9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SKY
    _no_line(line)


def _add_footer(slide):
    tag = _add_rounded_rect(slide, 0.55, 7.05, 2.8, 0.32, OCEAN)
    tag.text_frame.paragraphs[0].text = FOOTER_TEXT
    tag.text_frame.paragraphs[0].font.size = Pt(9)
    tag.text_frame.paragraphs[0].font.color.rgb = WHITE
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    deco = slide.shapes.add_textbox(Inches(8.5), Inches(7.05), Inches(1.0), Inches(0.3)).text_frame.paragraphs[0]
    deco.text = "🛒 📦 🌍"
    deco.font.size = Pt(11)
    deco.alignment = PP_ALIGN.RIGHT


def _add_tag(slide, x, y, text, bg=OCEAN, fg=WHITE, w=1.2):
    t = _add_rounded_rect(slide, x, y, w, 0.32, bg)
    p = t.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = fg
    p.alignment = PP_ALIGN.CENTER
    return t


# ── 各页模板 ──

def _add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_section_bg(slide)

    # 主标题卡片（白底亮色系）
    _add_shadow_card(slide, 0.9, 1.4, 8.2, 3.9)
    card = _add_rounded_rect(slide, 0.9, 1.4, 8.2, 3.9, WHITE, OCEAN, 2)

    _add_tag(slide, 1.2, 1.6, "Cross-Border E-Commerce", CORAL, WHITE, 2.6)

    p = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(7.6), Inches(1.1)).text_frame.paragraphs[0]
    p.text = "聊天记录质检工具"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY

    sp = slide.shapes.add_textbox(Inches(1.2), Inches(3.25), Inches(7.6), Inches(0.8)).text_frame.paragraphs[0]
    sp.text = "帮跨境电商销售团队解决「聊得多、检不过来、标准不统一」"
    sp.font.size = Pt(17)
    sp.font.color.rgb = OCEAN

    tp = slide.shapes.add_textbox(Inches(1.2), Inches(4.05), Inches(7.6), Inches(0.65)).text_frame.paragraphs[0]
    tp.text = "支持所有可导出聊天记录，或允许接入聊天记录 API 的平台"
    tp.font.size = Pt(13)
    tp.font.color.rgb = GRAY

    # 底部三个卖点标签
    for i, (icon, txt) in enumerate([("⚡", "10 分钟出日报"), ("🎯", "高风险自动标红"), ("📋", "结构化质检报告")]):
        x = 1.2 + i * 2.7
        chip = _add_rounded_rect(slide, x, 5.55, 2.4, 0.55, LIGHT_BLUE, OCEAN, 1)
        cp = chip.text_frame.paragraphs[0]
        cp.text = f"{icon}  {txt}"
        cp.font.size = Pt(12)
        cp.font.bold = True
        cp.font.color.rgb = OCEAN
        cp.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, title: str, subtitle: str = "", icon: str = "📌"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_section_bg(slide)

    _add_shadow_card(slide, 1.5, 2.3, 7.0, 2.5)
    card = _add_rounded_rect(slide, 1.5, 2.3, 7.0, 2.5, WHITE, OCEAN, 2)

    ip = slide.shapes.add_textbox(Inches(1.8), Inches(2.5), Inches(0.8), Inches(0.7)).text_frame.paragraphs[0]
    ip.text = icon
    ip.font.size = Pt(32)

    p = slide.shapes.add_textbox(Inches(2.5), Inches(2.55), Inches(5.8), Inches(0.9)).text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle:
        sp = slide.shapes.add_textbox(Inches(2.5), Inches(3.55), Inches(5.8), Inches(0.6)).text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(17)
        sp.font.color.rgb = OCEAN


def _add_content_slide(prs: Presentation, title: str, bullets: list[str], note: str = "", icon: str = "📋"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, title, icon)
    _add_footer(slide)

    card = _add_shadow_card(slide, 0.55, 1.2, 8.9, 5.5)
    tf = slide.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(8.2), Inches(4.9)).text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)

    if note:
        note_bar = _add_rounded_rect(slide, 0.55, 6.55, 8.9, 0.42, LIGHT_BLUE, SKY, 1)
        np = note_bar.text_frame.paragraphs[0]
        np.text = f"💡  {note}"
        np.font.size = Pt(11)
        np.font.color.rgb = OCEAN


def _add_pain_questions_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "做跨境电商销售管理，你是不是也遇到过这些？", "🤔")
    _add_footer(slide)

    questions = [
        ("💬", "每天上百通 WhatsApp / FB 聊天，主管根本看不过来，只能抽几通？"),
        ("⏰", "客户因回复慢、话术不对已流失，出问题才发现？"),
        ("📏", "不同主管点评口径不一样，客服不知道到底该怎么改？"),
        ("📝", "复盘全靠口头说，改进了什么、谁没改，没法追踪？"),
        ("🏆", "销冠怎么成交的，经验在脑子里，新人来了还是从头摸索？"),
    ]
    for i, (icon, q) in enumerate(questions):
        y = 1.2 + i * 1.1
        card = _add_shadow_card(slide, 0.55, y, 8.9, 0.92)
        ib = slide.shapes.add_textbox(Inches(0.75), Inches(y + 0.2), Inches(0.45), Inches(0.45)).text_frame.paragraphs[0]
        ib.text = icon
        ib.font.size = Pt(18)
        pb = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.22), Inches(7.9), Inches(0.55)).text_frame.paragraphs[0]
        pb.text = q
        pb.font.size = Pt(14)
        pb.font.color.rgb = DARK


def _add_pain_points_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "跨境电商销售管理的 5 大痛点", "⚠️")
    _add_footer(slide)

    pains = [
        ("01", "检不过来", "多平台对话量大，人工只能抽检 10–20%", CORAL),
        ("02", "发现太晚", "首响超时、话术失误，往往客户流失后才知道", CORAL),
        ("03", "标准不统一", "质检靠主管个人经验，不同人不同标准", GOLD),
        ("04", "反馈难落地", "口头点评没记录，改了什么、谁没改说不清", GOLD),
        ("05", "经验难复制", "成交话术留在销冠脑子里，扩团队、带新人都慢", OCEAN),
    ]
    for i, (num, title, desc, color) in enumerate(pains):
        y = 1.15 + i * 1.08
        card = _add_shadow_card(slide, 0.55, y, 8.9, 0.95)
        nb = _add_rounded_rect(slide, 0.75, y + 0.18, 0.45, 0.45, color)
        np = nb.text_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(11)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(Inches(1.35), Inches(y + 0.15), Inches(2.0), Inches(0.4)).text_frame.paragraphs[0]
        tb.text = title
        tb.font.size = Pt(15)
        tb.font.bold = True
        tb.font.color.rgb = NAVY

        db = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.18), Inches(5.8), Inches(0.55)).text_frame.paragraphs[0]
        db.text = desc
        db.font.size = Pt(13)
        db.font.color.rgb = GRAY


def _add_before_after_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "不用 vs 用了：差别在哪？", "⚖️")
    _add_footer(slide)

    for col, (label, bg, border, accent, items) in enumerate([
        ("✕  不用这套工具", RED_BG, RED_BORDER, RED_ACCENT, [
            "主管每天花 2–4 小时翻聊天，还是只能抽查",
            "低意向寒暄占大量时间，长对话没空看",
            "问题发现滞后，客户流失才意识到",
            "点评靠嘴说，没有记录，改没改说不清",
            "销冠经验带不走，新人试错成本高",
        ]),
        ("✓  用了这套工具", GREEN_BG, GREEN_BORDER, GREEN_ACCENT, [
            "上传 Excel，10–15 分钟出全量日报",
            "自动跳过低意向短对话，聚焦高价值会话",
            "不合格 / 高风险 / 红线自动标红，当天纠偏",
            "每通输出：问题 + 建议 + 下一步动作",
            "成交案例沉淀，标准越用越贴合业务",
        ]),
    ]):
        x = 0.55 + col * 4.55
        hdr = _add_rounded_rect(slide, x, 1.15, 4.25, 0.55, accent)
        hp = hdr.text_frame.paragraphs[0]
        hp.text = label
        hp.font.size = Pt(15)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER

        body = _add_rounded_rect(slide, x, 1.75, 4.25, 4.55, bg, border, 2)
        tf = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.95), Inches(3.85), Inches(4.2)).text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.space_after = Pt(9)


def _add_efficiency_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "效率能提升多少？（每天 200 通聊天）", "📈")
    _add_footer(slide)

    metrics = [
        ("15分钟", "质检耗时", "2~4h无效抽检 → 5+10min", "把最容易成交的客户全部握在手里", CORAL),
        ("5~10×", "覆盖范围", "抽检 → 全量", "不再靠抽样赌运气", OCEAN),
        ("60%↓", "无效质检", "跳过低意向短对话", "精力集中在高价值客户", TEAL),
        ("结构化", "每通有结论", "问题+建议+动作", "改法直接给到一线", GOLD),
        ("当天", "出报告", "不用等", "当日发现、当日纠偏", OCEAN),
        ("100%", "标准统一", "因人而异 → SOP", "公平考核、可对比", TEAL),
    ]

    for i, (num, label, change, note, color) in enumerate(metrics):
        col, row = i % 3, i // 3
        x, y = 0.55 + col * 3.1, 1.2 + row * 2.85

        card = _add_shadow_card(slide, x, y, 2.85, 2.55)
        top = _add_rounded_rect(slide, x, y, 2.85, 0.12, color)
        _no_line(top)

        np = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.25), Inches(2.55), Inches(0.55)).text_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(26)
        np.font.bold = True
        np.font.color.rgb = color

        lp = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.8), Inches(2.55), Inches(0.35)).text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(14)
        lp.font.bold = True
        lp.font.color.rgb = NAVY

        cp = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.15), Inches(2.55), Inches(0.5)).text_frame
        cp.word_wrap = True
        cp.paragraphs[0].text = change
        cp.paragraphs[0].font.size = Pt(10)
        cp.paragraphs[0].font.bold = True
        cp.paragraphs[0].font.color.rgb = GREEN_ACCENT

        dp = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.65), Inches(2.55), Inches(0.85)).text_frame
        dp.word_wrap = True
        dp.paragraphs[0].text = note
        dp.paragraphs[0].font.size = Pt(10)
        dp.paragraphs[0].font.color.rgb = GRAY


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], icon: str = "📊"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, title, icon)
    _add_footer(slide)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.55), Inches(1.15), Inches(8.9), Inches(0.48 * n_rows)
    )
    table = table_shape.table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = OCEAN
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LIGHT_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 表格外框
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.48 * n_rows + 0.15)
    )
    frame.fill.background()
    frame.line.color.rgb = CARD_BORDER
    frame.line.width = Pt(2)


def _add_table_slide_compact(prs, title, headers, rows, icon="📊", font_body=9, row_h=0.38):
    """行数较多的紧凑表格。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, title, icon)
    _add_footer(slide)

    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(
        n_rows, len(headers), Inches(0.5), Inches(1.12), Inches(9.0), Inches(row_h * n_rows)
    )
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = OCEAN
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_body + 1)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LIGHT_BLUE
            if i == len(rows) and row[0].startswith("合计"):
                cell.fill.fore_color.rgb = LIGHT_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_body)
                p.font.bold = row[0].startswith("合计") or j == 0
                p.font.color.rgb = DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_problems_impact_slide(prs: Presentation):
    """没有工具：销售工作中的问题与带来的影响。"""
    _add_table_slide_compact(
        prs,
        "没有这套工具：销售工作中会遇到的问题与影响",
        ["问题", "日常表现", "带来的影响"],
        [
            ["聊天检不过来", "每天上百通，主管只能抽检10~20%", "大量问题会话被漏掉，质量靠运气"],
            ["问题发现太晚", "首响超时、话术失误事后才知道", "客户已流失、订单已丢，补救成本高"],
            ["质检标准不统一", "不同主管点评口径不一", "客服无所适从，团队水平参差不齐"],
            ["反馈无法落地", "复盘靠口头，无结构化记录", "改了什么、谁没改说不清，同类错误反复犯"],
            ["高价值客户被淹没", "精力耗在低意向寒暄上", "真正该跟进的成交客户反而没空看"],
            ["销冠经验难复制", "成交路径留在个人脑子里", "扩团队慢，新人试错周期长、带教成本高"],
            ["管理半径受限", "1名主管精力有限，只能盯住约10人", "团队想扩到30人必须加主管编制"],
            ["数据无法量化", "没有合格率、风险率等统一指标", "绩效考核缺依据，管理决策靠感觉"],
        ],
        icon="⚠️",
        font_body=8,
        row_h=0.36,
    )


def _add_supervisor_scale_slide(prs: Presentation):
    """主管带人规模：10人 vs 30人。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "核心对比：1 名主管能带多少人？", "👥")
    _add_footer(slide)

    for col, (title, num, sub, items, bg, border, accent) in enumerate([
        ("没有系统", "10 人", "管理半径上限", [
            "主管每天2~4h翻聊天抽检",
            "只能覆盖10~20%会话",
            "30人团队需要 3 名主管",
        ], RED_BG, RED_BORDER, RED_ACCENT),
        ("拥有系统", "30 人", "同等精力管3倍团队", [
            "5分钟定位+10分钟分析",
            "全量过检，高风险自动标红",
            "30人团队只需 1 名主管",
        ], GREEN_BG, GREEN_BORDER, GREEN_ACCENT),
    ]):
        x = 0.55 + col * 4.55
        card = _add_shadow_card(slide, x, 1.15, 4.25, 3.8)
        _add_rounded_rect(slide, x, 1.15, 4.25, 0.55, accent)
        hp = slide.shapes.add_textbox(Inches(x), Inches(1.22), Inches(4.25), Inches(0.45)).text_frame.paragraphs[0]
        hp.text = title
        hp.font.size = Pt(15)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER

        np = slide.shapes.add_textbox(Inches(x), Inches(1.9), Inches(4.25), Inches(0.9)).text_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(42)
        np.font.bold = True
        np.font.color.rgb = accent
        np.alignment = PP_ALIGN.CENTER

        sp = slide.shapes.add_textbox(Inches(x), Inches(2.75), Inches(4.25), Inches(0.4)).text_frame.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(13)
        sp.font.color.rgb = GRAY
        sp.alignment = PP_ALIGN.CENTER

        tf = slide.shapes.add_textbox(Inches(x + 0.25), Inches(3.2), Inches(3.75), Inches(1.6)).text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)

    banner = _add_rounded_rect(slide, 0.55, 5.15, 8.9, 0.75, LIGHT_BLUE, OCEAN, 2)
    bp = banner.text_frame.paragraphs[0]
    bp.text = "人力账目：30人团队少配 2 名主管  →  每月节省约 ¥16,000 ~ ¥24,000（按主管月薪 ¥8,000~12,000 计）"
    bp.font.size = Pt(13)
    bp.font.bold = True
    bp.font.color.rgb = NAVY
    bp.alignment = PP_ALIGN.CENTER

    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.05), Inches(8.9), Inches(0.4)).text_frame.paragraphs[0]
    note.text = "※ 以下为参考估算，可按贵司实际薪资与团队规模替换数字"
    note.font.size = Pt(9)
    note.font.color.rgb = MUTED
    note.alignment = PP_ALIGN.CENTER


def _add_human_to_finance_slide(prs: Presentation):
    """人力成本转化为财务账目对比。"""
    _add_table_slide_compact(
        prs,
        "人力成本 → 财务账目对比（30人客服团队 · 月估算）",
        ["成本项", "没有系统", "拥有系统", "每月差额"],
        [
            ["主管质检工时", "2~4h/天×22天≈60h", "约2.5h/天≈12h", "节省约 ¥2,400"],
            ["主管编制（带人）", "3名主管", "1名主管", "节省 ¥16,000~24,000"],
            ["漏检致流失/客诉", "难以及时发现", "高风险当日标红", "减少损失 ¥5,000+"],
            ["新人带教试错", "2~4周上手", "报告当教材缩短周期", "节省 ¥3,000~5,000"],
            ["重复犯错无效工时", "无记录、同类问题反复", "结构化改进建议", "节省 ¥1,500~2,500"],
            ["合计（可量化部分）", "—", "—", "约 ¥28,000~39,000/月"],
        ],
        icon="💼",
        font_body=8,
        row_h=0.36,
    )


def _add_value_vs_cost_slide(prs: Presentation):
    """系统使用价值 vs 问题成本总和。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "系统价值：对冲「不用工具」的成本总和", "⚖️")
    _add_footer(slide)

    left = _add_rounded_rect(slide, 0.55, 1.15, 4.25, 4.85, RED_BG, RED_BORDER, 2)
    lp = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(3.85), Inches(0.4)).text_frame.paragraphs[0]
    lp.text = "不用工具 · 问题成本总和（月）"
    lp.font.size = Pt(13)
    lp.font.bold = True
    lp.font.color.rgb = RED_ACCENT

    costs = [
        ("主管多配编制", "¥16,000 ~ 24,000"),
        ("主管质检工时浪费", "¥2,000 ~ 4,000"),
        ("漏检客户流失", "¥5,000 +"),
        ("新人试错与带教", "¥3,000 ~ 5,000"),
        ("重复犯错无效工时", "¥1,500 ~ 2,500"),
    ]
    tf = slide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(3.85), Inches(3.5)).text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(costs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{k}：{v}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    tp = tf.add_paragraph()
    tp.text = "合计约 ¥28,000 ~ 39,000+ / 月"
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = RED_ACCENT

    right = _add_rounded_rect(slide, 5.2, 1.15, 4.25, 4.85, GREEN_BG, GREEN_BORDER, 2)
    rp = slide.shapes.add_textbox(Inches(5.4), Inches(1.3), Inches(3.85), Inches(0.4)).text_frame.paragraphs[0]
    rp.text = "拥有系统 · 可兑现价值"
    rp.font.size = Pt(13)
    rp.font.bold = True
    rp.font.color.rgb = GREEN_ACCENT

    values = [
        ("主管产能", "1人带30人，少配2名主管"),
        ("质检效率", "5分钟定位+10分钟分析"),
        ("质量覆盖", "100%过检，高风险标红"),
        ("管理闭环", "问题+建议+动作可追踪"),
        ("经验沉淀", "成交案例反哺标准"),
    ]
    rtf = slide.shapes.add_textbox(Inches(5.4), Inches(1.75), Inches(3.85), Inches(3.5)).text_frame
    rtf.word_wrap = True
    for i, (k, v) in enumerate(values):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"✓ {k}：{v}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    rtp = rtf.add_paragraph()
    rtp.text = "≈ 对冲上述成本总和，且管理半径扩大 3 倍"
    rtp.font.size = Pt(12)
    rtp.font.bold = True
    rtp.font.color.rgb = GREEN_ACCENT

    arrow = slide.shapes.add_textbox(Inches(4.55), Inches(3.2), Inches(0.9), Inches(0.5)).text_frame.paragraphs[0]
    arrow.text = "→"
    arrow.font.size = Pt(28)
    arrow.font.color.rgb = OCEAN
    arrow.alignment = PP_ALIGN.CENTER

    bottom = _add_rounded_rect(slide, 0.55, 6.15, 8.9, 0.55, OCEAN)
    bp = bottom.text_frame.paragraphs[0]
    bp.text = "结论：系统价值 ≈ 每月节省 2~3 万元人力与管理成本 + 减少不可量化的客户流失"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER


def _add_system_solves_slide(prs: Presentation):
    """拥有系统后能解决/提升什么。"""
    _add_table_slide_compact(
        prs,
        "拥有系统后：能解决什么、提升什么",
        ["原问题", "系统怎么做", "提升结果"],
        [
            ["检不过来", "全量自动分档质检", "100%覆盖，不再靠抽检"],
            ["发现太晚", "高风险/不合格标红", "当时定位、当时纠偏"],
            ["标准不统一", "统一SOP+AI判断", "所有人同一把尺子"],
            ["反馈难落地", "问题+建议+动作", "可直接下发、可追踪"],
            ["高价值客户淹没", "跳过低意向短对话", "精力集中在最易成交客户"],
            ["主管只能带10人", "15分钟完成质检分析", "1主管可带约30人"],
            ["经验难复制", "成交案例持续学习", "新人有案可查"],
            ["无法量化考核", "合格率/风险率数据", "绩效有据、公平对比"],
        ],
        icon="✅",
        font_body=8,
        row_h=0.36,
    )


def _add_detailed_workflow_slide(prs: Presentation):
    """详细运用流程。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "运用流程：从数据到管理动作", "🔄")
    _add_footer(slide)

    steps = [
        ("1", "准备数据", "SaleSmartly：API 自动拉取\n其他平台：手动上传 Excel"),
        ("2", "确认配置", "核对列映射、时间范围\n筛选要质检的会话"),
        ("3", "启动质检", "点击开始，等待十余分钟\n系统自动分档+AI分析"),
        ("4", "下载报告", "获取 Excel 双表报告\n优质客户 + 低意向客户"),
        ("5", "定位问题", "筛「重点档」+ 高风险标红\n5分钟找到该查谁"),
        ("6", "分析纠偏", "看问题/建议/跟进动作\n10分钟完成分析下发"),
        ("7", "复盘沉淀", "周会案例教学\n成交客户反哺标准"),
    ]
    start_x = 0.35
    step_w = 1.28
    gap = 0.08
    for i, (num, name, desc) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        card = _add_shadow_card(slide, x, 1.2, step_w, 4.5)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.39), Inches(1.4), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = OCEAN
        _no_line(circle)
        cn = circle.text_frame.paragraphs[0]
        cn.text = num
        cn.font.size = Pt(14)
        cn.font.bold = True
        cn.font.color.rgb = WHITE
        cn.alignment = PP_ALIGN.CENTER

        np = slide.shapes.add_textbox(Inches(x + 0.05), Inches(2.05), Inches(step_w - 0.1), Inches(0.4)).text_frame.paragraphs[0]
        np.text = name
        np.font.size = Pt(11)
        np.font.bold = True
        np.font.color.rgb = NAVY
        np.alignment = PP_ALIGN.CENTER

        dp = slide.shapes.add_textbox(Inches(x + 0.05), Inches(2.5), Inches(step_w - 0.1), Inches(2.8)).text_frame
        dp.word_wrap = True
        dp.paragraphs[0].text = desc
        dp.paragraphs[0].font.size = Pt(9)
        dp.paragraphs[0].font.color.rgb = GRAY
        dp.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            ap = slide.shapes.add_textbox(Inches(x + step_w), Inches(2.0), Inches(gap + 0.02), Inches(0.35)).text_frame.paragraphs[0]
            ap.text = "›"
            ap.font.size = Pt(14)
            ap.font.color.rgb = CORAL


def _add_video_placeholder_slide(prs: Presentation):
    """操作视频占位页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "操作演示视频", "🎬")
    _add_footer(slide)

    frame = _add_rounded_rect(slide, 0.55, 1.15, 8.9, 4.2, PLACEHOLDER_BG, PLACEHOLDER_BORDER, 2)
    frame.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    play = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.35), Inches(2.6), Inches(1.3), Inches(1.3))
    play.fill.solid()
    play.fill.fore_color.rgb = OCEAN
    _no_line(play)
    pp = play.text_frame.paragraphs[0]
    pp.text = "▶"
    pp.font.size = Pt(36)
    pp.font.color.rgb = WHITE
    pp.alignment = PP_ALIGN.CENTER

    lp = slide.shapes.add_textbox(Inches(0.55), Inches(4.1), Inches(8.9), Inches(0.5)).text_frame.paragraphs[0]
    lp.text = "【操作视频占位】建议录制：上传 Excel → 质检进度 → 下载报告 → 筛高风险"
    lp.font.size = Pt(14)
    lp.font.bold = True
    lp.font.color.rgb = OCEAN
    lp.alignment = PP_ALIGN.CENTER

    steps_card = _add_rounded_rect(slide, 0.55, 5.5, 8.9, 1.0, LIGHT_BLUE, SKY, 1)
    stf = steps_card.text_frame
    stf.word_wrap = True
    lines = [
        "视频建议涵盖：① 登录网页版  ② 选择数据来源（API/Excel）  ③ 上传并确认列映射",
        "④ 查看质检进度与统计  ⑤ 下载报告  ⑥ 筛选高风险并解读改法",
    ]
    for i, line in enumerate(lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = OCEAN


def _add_problem_solution_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "这个工具帮你解决什么？", "🎯")
    _add_footer(slide)

    rows = [
        ("你的痛点", "工具怎么解决", "你得到什么"),
        ("聊天太多，检不过来", "自动分档 + AI 批量质检", "10 分钟出全量日报"),
        ("不知道先查谁", "高风险 / 不合格自动标红", "打开就知道先处理谁"),
        ("标准因人而异", "按公司 SOP 统一判断", "所有人同一套尺子"),
        ("点评完就忘了", "结构化：问题+建议+动作", "可直接下发、可追踪"),
        ("短对话浪费精力", "低意向自动跳过", "精力集中在高价值客户"),
        ("销冠经验带不走", "成交案例持续学习", "新人有案可查"),
    ]

    table = slide.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.15), Inches(8.9), Inches(0.46 * len(rows))).table
    for j, w in enumerate([2.6, 3.3, 3.0]):
        table.columns[j].width = Inches(w)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if i > 0 else 12)
                p.font.bold = i == 0
                p.font.color.rgb = WHITE if i == 0 else DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_one_liner_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_section_bg(slide)

    _add_shadow_card(slide, 1.0, 1.8, 8.0, 3.5)
    _add_rounded_rect(slide, 1.0, 1.8, 8.0, 3.5, WHITE, OCEAN, 2)

    p = slide.shapes.add_textbox(Inches(1.4), Inches(2.2), Inches(7.2), Inches(1.4)).text_frame.paragraphs[0]
    p.text = "把聊天记录表格丢进去\n当天拿到「该查谁、该怎么改」的 Excel 报告"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    chips = ["不用逐通翻", "不用靠人记", "不用猜标准"]
    for i, txt in enumerate(chips):
        x = 1.6 + i * 2.3
        c = _add_rounded_rect(slide, x, 4.0, 2.0, 0.5, LIGHT_BLUE, OCEAN, 1)
        cp = c.text_frame.paragraphs[0]
        cp.text = f"✓ {txt}"
        cp.font.size = Pt(13)
        cp.font.bold = True
        cp.font.color.rgb = OCEAN
        cp.alignment = PP_ALIGN.CENTER


def _add_flow_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "怎么用？三步搞定", "🚀")
    _add_footer(slide)

    steps = [
        ("📤", "1", "导入", "上传当天聊天\nExcel 表格"),
        ("⏳", "2", "等待", "10~15 分钟\n自动质检"),
        ("📊", "3", "看报告", "筛高风险\n下发改进"),
    ]
    start_x = 0.9
    step_w = 2.7
    gap = 0.35

    for i, (emoji, num, name, desc) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        card = _add_shadow_card(slide, x, 1.5, step_w, 4.2)

        eb = slide.shapes.add_textbox(Inches(x + 0.95), Inches(1.7), Inches(0.8), Inches(0.6)).text_frame.paragraphs[0]
        eb.text = emoji
        eb.font.size = Pt(28)
        eb.alignment = PP_ALIGN.CENTER

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.05), Inches(2.35), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = OCEAN
        _no_line(circle)
        cn = circle.text_frame.paragraphs[0]
        cn.text = num
        cn.font.size = Pt(18)
        cn.font.bold = True
        cn.font.color.rgb = WHITE
        cn.alignment = PP_ALIGN.CENTER

        np = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.1), Inches(step_w - 0.4), Inches(0.45)).text_frame.paragraphs[0]
        np.text = name
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = NAVY
        np.alignment = PP_ALIGN.CENTER

        dp = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.6), Inches(step_w - 0.4), Inches(1.0)).text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(13)
        dp.font.color.rgb = GRAY
        dp.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            ap = slide.shapes.add_textbox(Inches(x + step_w), Inches(3.0), Inches(gap), Inches(0.5)).text_frame.paragraphs[0]
            ap.text = "→"
            ap.font.size = Pt(26)
            ap.font.color.rgb = CORAL
            ap.alignment = PP_ALIGN.CENTER


def _png_size(path: Path) -> tuple[int, int]:
    with path.open("rb") as f:
        f.seek(16)
        w, h = struct.unpack(">II", f.read(8))
    return w, h


def _add_picture_fit(slide, path: Path, x: float, y: float, max_w: float, max_h: float):
    """按比例缩放图片，放入指定区域并居中。"""
    iw, ih = _png_size(path)
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    ox = x + (max_w - w) / 2
    oy = y + (max_h - h) / 2
    frame = _add_rounded_rect(slide, x, y, max_w, max_h, WHITE, CARD_BORDER, 1.5)
    pic = slide.shapes.add_picture(str(path), Inches(ox), Inches(oy), width=Inches(w))
    return frame, pic


def _add_screenshot_slide(prs, title, image_name: str, caption, side_note=""):
    """嵌入真实截图（ppt_assets 目录下）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, title, "🖼️")
    _add_footer(slide)

    image_path = ASSETS_DIR / image_name
    img_x, img_y = 0.55, 1.15
    img_w = 7.5 if side_note else 8.9
    img_h = 4.7

    if image_path.is_file():
        _add_picture_fit(slide, image_path, img_x, img_y, img_w, img_h)
    else:
        ph = _add_rounded_rect(slide, img_x, img_y, img_w, img_h, PLACEHOLDER_BG, PLACEHOLDER_BORDER, 2)
        ph.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        lp = slide.shapes.add_textbox(Inches(img_x), Inches(img_y + 2.0), Inches(img_w), Inches(0.5)).text_frame.paragraphs[0]
        lp.text = f"缺少截图：{image_name}"
        lp.font.size = Pt(14)
        lp.font.color.rgb = GRAY
        lp.alignment = PP_ALIGN.CENTER

    cap_card = _add_rounded_rect(slide, 0.55, 5.95, 8.9 if not side_note else 6.0, 0.48, LIGHT_BLUE, SKY, 1)
    cp = cap_card.text_frame.paragraphs[0]
    cp.text = f"💡  {caption}"
    cp.font.size = Pt(11)
    cp.font.color.rgb = OCEAN

    if side_note:
        _add_shadow_card(slide, 6.65, 1.15, 2.8, 4.7)
        ntp = slide.shapes.add_textbox(Inches(6.85), Inches(1.35), Inches(2.4), Inches(0.4)).text_frame.paragraphs[0]
        ntp.text = "看报告时关注"
        ntp.font.size = Pt(12)
        ntp.font.bold = True
        ntp.font.color.rgb = CORAL
        ntf = slide.shapes.add_textbox(Inches(6.85), Inches(1.8), Inches(2.4), Inches(3.6)).text_frame
        ntf.word_wrap = True
        for i, line in enumerate(side_note.split("\n")):
            p = ntf.paragraphs[0] if i == 0 else ntf.add_paragraph()
            p.text = f"▸ {line}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)


def _add_daily_use_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "用了之后，日常管理会变成这样", "📅")
    _add_footer(slide)

    uses = [
        ("⚡ 需要质检时", "上传聊天记录 → 十余分钟出报告", "筛高风险，当时纠偏", OCEAN),
        ("📆 每周复盘会", "用「重点档」长对话做案例", "不用现场翻记录", TEAL),
        ("👩‍🏫 新人培训", "拿改进建议 + 跟进动作当教材", "有具体话术参考", GOLD),
        ("📈 绩效考核", "对比各客服合格率、风险率", "有数据、更公平", CORAL),
    ]
    for i, (when, action, result, color) in enumerate(uses):
        y = 1.15 + i * 1.38
        card = _add_shadow_card(slide, 0.55, y, 8.9, 1.2)
        bar = _add_rounded_rect(slide, 0.55, y, 0.1, 1.2, color)
        _no_line(bar)

        wb = slide.shapes.add_textbox(Inches(0.85), Inches(y + 0.15), Inches(1.8), Inches(0.4)).text_frame.paragraphs[0]
        wb.text = when
        wb.font.size = Pt(13)
        wb.font.bold = True
        wb.font.color.rgb = color

        ab = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.2), Inches(3.5), Inches(0.8)).text_frame.paragraphs[0]
        ab.text = action
        ab.font.size = Pt(12)
        ab.font.color.rgb = DARK

        rb = _add_rounded_rect(slide, 6.3, y + 0.25, 2.9, 0.65, LIGHT_BLUE, SKY, 1)
        rp = rb.text_frame.paragraphs[0]
        rp.text = f"→ {result}"
        rp.font.size = Pt(11)
        rp.font.bold = True
        rp.font.color.rgb = GREEN_ACCENT


def _add_experience_and_upload_slide(prs: Presentation):
    """结尾强调：成交客户学习 + 经验注入 + 数据接入方式。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "不止质检：让成交经验持续注入系统", "🔄")
    _add_footer(slide)

    # 左：成交学习闭环
    _add_shadow_card(slide, 0.55, 1.15, 4.35, 5.35)
    _add_rounded_rect(slide, 0.55, 1.15, 4.35, 0.5, TEAL)
    lt = slide.shapes.add_textbox(Inches(0.55), Inches(1.22), Inches(4.35), Inches(0.4)).text_frame.paragraphs[0]
    lt.text = "成交客户学习 · 经验注入"
    lt.font.size = Pt(14)
    lt.font.bold = True
    lt.font.color.rgb = WHITE
    lt.alignment = PP_ALIGN.CENTER

    loop_steps = [
        "自动收录已成交客户的完整聊天记录",
        "AI 分析成交路径、话术与心理轨迹",
        "提炼顾虑化解、逼单跟单等最佳实践",
        "沉淀为可复用知识，反哺质检标准",
        "质检越用越准，新人培训有真实案例",
    ]
    ltf = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(3.85), Inches(4.4)).text_frame
    ltf.word_wrap = True
    for i, step in enumerate(loop_steps):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"{'①②③④⑤'[i]}  {step}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)

    highlight = _add_rounded_rect(slide, 0.7, 5.55, 4.05, 0.75, LIGHT_BLUE, TEAL, 1.5)
    hp = highlight.text_frame.paragraphs[0]
    hp.text = "从「个人经验」到「组织资产」\n销冠走了，成交方法还在系统里"
    hp.font.size = Pt(11)
    hp.font.bold = True
    hp.font.color.rgb = OCEAN
    hp.alignment = PP_ALIGN.CENTER

    # 右：数据接入
    _add_shadow_card(slide, 5.1, 1.15, 4.35, 5.35)
    _add_rounded_rect(slide, 5.1, 1.15, 4.35, 0.5, OCEAN)
    rt = slide.shapes.add_textbox(Inches(5.1), Inches(1.22), Inches(4.35), Inches(0.4)).text_frame.paragraphs[0]
    rt.text = "聊天记录怎么进来？"
    rt.font.size = Pt(14)
    rt.font.bold = True
    rt.font.color.rgb = WHITE
    rt.alignment = PP_ALIGN.CENTER

    # SaleSmartly 自动
    sm_card = _add_rounded_rect(slide, 5.35, 1.85, 3.85, 1.55, GREEN_BG, GREEN_BORDER, 2)
    sp = slide.shapes.add_textbox(Inches(5.5), Inches(1.95), Inches(3.55), Inches(1.35)).text_frame
    sp.word_wrap = True
    p0 = sp.paragraphs[0]
    p0.text = "SaleSmartly  ·  自动接入"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = GREEN_ACCENT
    p1 = sp.add_paragraph()
    p1.text = "对接 API，聊天记录自动拉取入库\n质检、成交学习均可定时自动跑"
    p1.font.size = Pt(11)
    p1.font.color.rgb = DARK
    p1.space_before = Pt(6)

    # 其他平台手动
    mn_card = _add_rounded_rect(slide, 5.35, 3.55, 3.85, 1.55, LIGHT_BLUE, SKY, 2)
    mp = slide.shapes.add_textbox(Inches(5.5), Inches(3.65), Inches(3.55), Inches(1.35)).text_frame
    mp.word_wrap = True
    m0 = mp.paragraphs[0]
    m0.text = "其他平台  ·  手动上传"
    m0.font.size = Pt(13)
    m0.font.bold = True
    m0.font.color.rgb = OCEAN
    m1 = mp.add_paragraph()
    m1.text = "凡能导出聊天记录 Excel\n或通过 API 接入的平台均支持"
    m1.font.size = Pt(11)
    m1.font.color.rgb = DARK
    m1.space_before = Pt(6)

    note = _add_rounded_rect(slide, 5.35, 5.25, 3.85, 0.95, WHITE, CARD_BORDER, 1)
    np = note.text_frame.paragraphs[0]
    np.text = "质检与成交学习共用同一套数据\n无论自动还是手动，进系统后流程一致"
    np.font.size = Pt(10)
    np.font.color.rgb = GRAY
    np.alignment = PP_ALIGN.CENTER

    # 底部总结条
    banner = _add_rounded_rect(slide, 0.55, 6.55, 8.9, 0.42, CORAL)
    bp = banner.text_frame.paragraphs[0]
    bp.text = "质检发现问题 → 成交案例沉淀经验 → 经验注入标准 → 团队越用越强"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER


def _add_ai_agent_learning_slide(prs: Presentation):
    """重点说明：AI 智能体如何持续学习。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_light_bg(slide)
    _add_header(slide, "【重点】怎么学习？打造专属于业务的 AI 智能体", "🧠")
    _add_footer(slide)

    _add_tag(slide, 0.55, 1.05, "核心能力", CORAL, WHITE, 1.1)

    core = _add_rounded_rect(slide, 0.55, 1.12, 8.9, 0.95, WHITE, OCEAN, 2)
    cp = core.text_frame.paragraphs[0]
    cp.text = (
        "通过 AI 大模型，创建「专属于业务的智能体」——"
        "不断加深对客户的了解、对业务的了解，越用越懂你的生意。"
    )
    cp.font.size = Pt(13)
    cp.font.bold = True
    cp.font.color.rgb = NAVY
    cp.alignment = PP_ALIGN.CENTER

    learn_steps = [
        ("📥", "吸收", "成交聊天记录\n质检案例与业务文档", OCEAN),
        ("🔍", "理解", "AI 分析心理轨迹\n顾虑化解与成交路径", TEAL),
        ("💉", "注入", "经验写入智能体\n反哺质检与话术标准", CORAL),
    ]
    for i, (icon, name, desc, color) in enumerate(learn_steps):
        x = 0.7 + i * 3.05
        _add_shadow_card(slide, x, 2.2, 2.75, 2.1)
        _add_rounded_rect(slide, x, 2.2, 2.75, 0.12, color)

        ib = slide.shapes.add_textbox(Inches(x + 0.95), Inches(2.35), Inches(0.85), Inches(0.5)).text_frame.paragraphs[0]
        ib.text = icon
        ib.font.size = Pt(22)
        ib.alignment = PP_ALIGN.CENTER

        nb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.9), Inches(2.45), Inches(0.35)).text_frame.paragraphs[0]
        nb.text = name
        nb.font.size = Pt(15)
        nb.font.bold = True
        nb.font.color.rgb = color
        nb.alignment = PP_ALIGN.CENTER

        db = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.3), Inches(2.45), Inches(0.85)).text_frame
        db.word_wrap = True
        db.paragraphs[0].text = desc
        db.paragraphs[0].font.size = Pt(10)
        db.paragraphs[0].font.color.rgb = GRAY
        db.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < 2:
            ar = slide.shapes.add_textbox(Inches(x + 2.75), Inches(2.95), Inches(0.3), Inches(0.4)).text_frame.paragraphs[0]
            ar.text = "→"
            ar.font.size = Pt(18)
            ar.font.color.rgb = SKY

    loop_card = _add_rounded_rect(slide, 0.55, 4.45, 8.9, 0.62, LIGHT_BLUE, TEAL, 1.5)
    lp = loop_card.text_frame.paragraphs[0]
    lp.text = "智能体持续学习：每多一笔成交、每多一次复盘，对业务与客户的理解就加深一层"
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = OCEAN
    lp.alignment = PP_ALIGN.CENTER

    adapt_title = slide.shapes.add_textbox(Inches(0.55), Inches(5.2), Inches(8.9), Inches(0.32)).text_frame.paragraphs[0]
    adapt_title.text = "能适应不同地区 · 不同风格 · 不同时期的客户"
    adapt_title.font.size = Pt(13)
    adapt_title.font.bold = True
    adapt_title.font.color.rgb = NAVY
    adapt_title.alignment = PP_ALIGN.CENTER

    adapts = [
        ("🌍 不同地区", "欧美/东南亚/中东\n沟通习惯各不相同", OCEAN),
        ("🎭 不同风格", "理性/犹豫/冲动型\n对症话术各不相同", TEAL),
        ("📅 不同时期", "旺季促销/淡季维护\n客户心态随周期变化", GOLD),
    ]
    for i, (title, desc, color) in enumerate(adapts):
        x = 0.55 + i * 3.05
        _add_shadow_card(slide, x, 5.55, 2.85, 1.35)
        _add_rounded_rect(slide, x, 5.55, 2.85, 0.1, color)

        at = slide.shapes.add_textbox(Inches(x + 0.12), Inches(5.68), Inches(2.6), Inches(0.32)).text_frame.paragraphs[0]
        at.text = title
        at.font.size = Pt(11)
        at.font.bold = True
        at.font.color.rgb = color

        ad = slide.shapes.add_textbox(Inches(x + 0.12), Inches(6.0), Inches(2.6), Inches(0.75)).text_frame
        ad.word_wrap = True
        ad.paragraphs[0].text = desc
        ad.paragraphs[0].font.size = Pt(9)
        ad.paragraphs[0].font.color.rgb = GRAY


def _add_closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_section_bg(slide)

    _add_shadow_card(slide, 0.9, 1.5, 8.2, 4.5)
    _add_rounded_rect(slide, 0.9, 1.5, 8.2, 4.5, WHITE, OCEAN, 2)

    p = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(7.4), Inches(0.7)).text_frame.paragraphs[0]
    p.text = "一句话总结"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    for i, (icon, txt, color) in enumerate([
        ("✕", "不用：靠人翻、靠抽查、靠记忆 → 问题发现慢、标准不统一", RED_ACCENT),
        ("✓", "用了：当天全量出报告，高风险标红，改法直接给到一线", GREEN_ACCENT),
    ]):
        _add_rounded_rect(slide, 1.3, 2.7 + i * 1.1, 7.4, 0.85,
                          RED_BG if i == 0 else GREEN_BG,
                          RED_BORDER if i == 0 else GREEN_BORDER, 1.5)
        rp = slide.shapes.add_textbox(Inches(1.5), Inches(2.9 + i * 1.1), Inches(7.0), Inches(0.6)).text_frame.paragraphs[0]
        rp.text = f"{icon}  {txt}"
        rp.font.size = Pt(14)
        rp.font.color.rgb = color

    cp = slide.shapes.add_textbox(Inches(1.3), Inches(5.0), Inches(7.4), Inches(0.5)).text_frame.paragraphs[0]
    cp.text = "欢迎拿聊天记录试用：SaleSmartly 可自动接入，其他平台手动上传即可"
    cp.font.size = Pt(15)
    cp.font.color.rgb = OCEAN
    cp.alignment = PP_ALIGN.CENTER

    qp = slide.shapes.add_textbox(Inches(1.3), Inches(5.7), Inches(7.4), Inches(0.5)).text_frame.paragraphs[0]
    qp.text = "Q & A"
    qp.font.size = Pt(22)
    qp.font.color.rgb = CORAL
    qp.alignment = PP_ALIGN.CENTER


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)
    _add_pain_questions_slide(prs)
    _add_problems_impact_slide(prs)
    _add_pain_points_slide(prs)

    _add_section_slide(prs, "成本与价值", "人力账目 · 财务对比", "💼")
    _add_supervisor_scale_slide(prs)
    _add_human_to_finance_slide(prs)
    _add_value_vs_cost_slide(prs)
    _add_system_solves_slide(prs)

    _add_section_slide(prs, "不用 vs 用了", "差别一目了然", "⚖️")
    _add_before_after_slide(prs)
    _add_efficiency_slide(prs)

    _add_table_slide(
        prs, "数据对比一览（每天 200 通聊天）",
        ["对比项", "不用", "用了", "提升"],
        [
            ["质检耗时", "2~4 小时无效抽检", "5分钟定位+10分钟分析", "掌握最易成交客户"],
            ["会话覆盖率", "10~20% 抽检", "100% 过检", "5~10 倍覆盖"],
            ["问题发现", "事后才知道", "当天标红预警", "当日响应"],
            ["反馈形式", "口头点评", "Excel 结构化", "可追踪"],
            ["质检标准", "因人而异", "统一 SOP", "公平可对比"],
        ],
    )

    _add_section_slide(prs, "工具能帮你做什么", "痛点 → 方案 → 结果", "🎯")
    _add_problem_solution_slide(prs)
    _add_one_liner_slide(prs)

    _add_section_slide(prs, "操作流程", "运用流程 · 操作演示", "🎬")
    _add_detailed_workflow_slide(prs)
    _add_video_placeholder_slide(prs)

    _add_section_slide(prs, "报告长什么样", "真实界面截图", "🖼️")

    _add_screenshot_slide(
        prs, "上传表格，等几分钟", "01_upload.png",
        "网页版支持 API 拉取或上传 Excel，不需要懂技术。",
    )
    _add_screenshot_slide(
        prs, "报告：一眼看到该查谁", "02_report.png",
        "不合格、高风险自动标红，客户阶段与主要顾虑一目了然。",
        side_note="先筛「重点档」\n再看标红项\n直接用于复盘",
    )
    _add_screenshot_slide(
        prs, "每通对话都有改法", "03_actions.png",
        "不只是说「不合格」，还告诉客服具体该怎么改、下一步做什么。",
    )

    _add_daily_use_slide(prs)

    _add_content_slide(
        prs, "适合哪些跨境电商团队？",
        [
            "销售 / 客服主管：每天需要盯 WhatsApp、FB、IG 等多渠道聊天质量",
            "对话量大、渠道多的团队：标准难统一，靠人工抽检覆盖不全",
            "正在扩团队的业务：想把销冠经验沉淀下来，缩短新人上手周期",
            "已用 SalesMartly 的团队：可直接对接聊天数据，减少手工导出",
        ],
        note="提供桌面版（本机）和网页版（浏览器），可按团队情况选择。",
        icon="🌍",
    )

    _add_section_slide(prs, "长期价值", "下一页重点：AI 智能体如何学习", "🌱")
    _add_ai_agent_learning_slide(prs)
    _add_experience_and_upload_slide(prs)

    _add_closing_slide(prs)
    return prs


def main():
    prs = build_presentation()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"已生成：{OUTPUT}")


if __name__ == "__main__":
    main()
